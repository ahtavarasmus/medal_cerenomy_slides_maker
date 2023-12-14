import sys
import tempfile
from PyQt6.QtWidgets import QApplication, QWidget, QVBoxLayout, QPushButton, QFileDialog, QLabel, QLineEdit
from PyQt6.QtGui import QMovie
from PyQt6.QtCore import pyqtSlot, Qt, QThread, pyqtSignal
from time import sleep
from pptx import Presentation
from pptx.util import Pt
from PIL import Image
import csv,sys,logging
from pptx.util import Inches
from datetime import datetime

import imageio

import aspose.slides as slides
import os


logging.basicConfig(level=logging.DEBUG,  # Set the lowest level to capture everything
                    format='%(asctime)s - %(levelname)s - %(message)s',
                    filename='app_log.txt',  # File where logs will be written
                    filemode='w')  # 'w' for overwrite each time; use 'a' to append


basedir = os.path.dirname(__file__)

    # Define GIF paths and output path
gif_paths_same_height = {"fin":"same_height\\Flag_of_Finland1.gif",
            "swe":"same_height\\Flag_of_Sweden1.gif",
            "nor":"same_height\\Flag_of_Norway1.gif",
            "est":"same_height\\Flag_of_Estonia1.gif",
            "isl":"same_height\\Flag_of_Iceland1.gif",
            "ger":"same_height\\Flag_of_Germany1.gif",
            "ukr":"same_height\\Flag_of_Ukraine1.gif",
            "usa":"same_height\\Flag_of_the_United_States1.gif",
            }


# Define GIF paths and output path
gif_paths_same_width = {"fin":"same_width\\Flag_of_Finland.gif",
            "swe":"same_width\\Flag_of_Sweden.gif",
            "nor":"same_width\\Flag_of_Norway.gif",
            "est":"same_width\\Flag_of_Estonia.gif",
            "isl":"same_width\\Flag_of_Iceland.gif",
            "ger":"same_width\\Flag_of_Germany.gif",
            "ukr":"same_width\\Flag_of_Ukraine.gif",
            "usa":"same_width\\Flag_of_the_United_States.gif",
            }



def read_csv(filename):
    logging.info(f"Reading {filename}")
    try:
        with open(filename, mode='r', encoding='iso-8859-1') as file:
            data = dict()
            order = []
            reader = csv.reader(file, delimiter=';')
            for row in reader:
                data[row[0]] = [row[1:]] if row[0] not in data.keys() else data[row[0]] + [row[1:]]

            for key in data.keys():
                order.append(key)
            return data,order
    except FileNotFoundError:
        logging.error(f"File {filename} not found")
        sys.exit(1)
    logging.info(f"Read {len(data)} rows")
    logging.info(f"Sarjat: {data.keys()}")


def combine_gifs(gif_paths, sarja):
    output_path = os.path.join(os.getcwd(),"Combined_Flag_Podium_{sarja}.gif")
    logging.info("Combining GIFs")
    assert len(gif_paths) == 3, f"only three GIFs must be provided, got {len(gif_paths)} with {gif_paths}"

    # Create a temporary file for the combined gif
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.gif')
    output_path = temp_file.name
    temp_file.close()  # Close the file so that other functions can use iti

    # Read individual GIFs and get their frames
    gif_frames = [imageio.mimread(gif) for gif in gif_paths]
    extra_space = 50  # Additional horizontal space between flags

    # Get dimensions for combined gif
    combined_width = sum(Image.open(gif_paths[i]).size[0] for i in range(len(gif_paths))) + 150
    combined_height = max(Image.open(gif_paths[i]).size[1] for i in range(len(gif_paths))) + 300

    # Create frames for combined gif
    combined_frames = []
    for i in range(min(len(frames) for frames in gif_frames)):  # Loop through each frame
        # Create a new image with white background
        with Image.new('RGBA', (combined_width, combined_height), 'WHITE') as new_frame:

            x_offset = 20
            for idx, frames in enumerate(gif_frames):
                # Paste the individual frame
                frame_image = Image.fromarray(frames[i])

                y_offset = 150  # Default y-offset
                if idx == 1:  # If it's the middle (winner) flag
                    y_offset = 20  # Move it higher up
                elif idx == 2:
                    y_offset = 250

                new_frame.paste(frame_image, (x_offset, y_offset))
                x_offset += frame_image.width + extra_space  # Add extra space after each flag

            combined_frames.append(new_frame)

    # Save frames as new GIF
    imageio.mimsave(output_path, [frame for frame in combined_frames], format='GIF', duration=0.2, loop=0)
    logging.info(f"Saved combined gif to {output_path}")
    return output_path


basedir = os.path.dirname(__file__)

def abs_path(relative_path):
    """ Get the absolute path to the resource """

    return os.path.join(basedir, relative_path)

def create_presentation(day_number, flag_size, number_of_slides,prefix,suffix,data):
    logging.info("Creating presentation")
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]

    winners_acronyms = []
    for sarja in data.keys():
        # creating 100 slides for each sarja to create animation
        # sorting based on tulos
        logging.info(f"sarja {sarja}, data[sarja]: ")
        for row in data[sarja]:
            logging.info(row)
        #logging.info(f"Sorting sarja {sarja}, data[sarja]: ")
        #sorted(data[sarja], key=lambda x: x[4])

        if len(data[sarja]) < 3:
            logging.error(f"not enough rows for sarja {sarja}")
            continue
        first_country = data[sarja][0][3].lower()
        second_country = data[sarja][1][3].lower()
        third_country = data[sarja][2][3].lower()
        if flag_size == "height":
            if first_country in gif_paths_same_height.keys() and second_country in gif_paths_same_height.keys() and third_country in gif_paths_same_height.keys():
                first = gif_paths_same_height[first_country]
                second = gif_paths_same_height[second_country]
                third = gif_paths_same_height[third_country]
            else:
                logging.error(f"didn't find country for sarja {sarja} with countries {first_country}, {second_country}, {third_country}")
                continue
        elif flag_size == "width":
            if first_country in gif_paths_same_width.keys() and second_country in gif_paths_same_width.keys() and third_country in gif_paths_same_width.keys():
                first = gif_paths_same_width[first_country]
                second = gif_paths_same_width[second_country]
                third = gif_paths_same_width[third_country]
            else:
                logging.error(f"didn't find country for sarja {sarja} with countries {first_country}, {second_country}, {third_country}")
                continue
        else:
            logging.error(f"unknown flag size {flag_size}")
            exit(1)
        if not first or not second or not third:
            logging.error(f"didn't find gif for sarja {sarja} with countries {first_country}, {second_country}, {third_country}")
            continue
        logging.info(f"Sarja: {sarja}, 1: {first}, 2: {second}, 3: {third}")

        combined_gif_path = combine_gifs([second, first, third], sarja)
        winners_acronyms.append(first_country.upper())
        logging.info(f"Making slide for showing sarja {sarja}")

        slide = prs.slides.add_slide(slide_layout)
        # adding stationary large title text to the middle
        title = slide.shapes.title
        title.text = prefix+sarja+suffix
        title.text_frame.paragraphs[0].font.size = Pt(44)  # Set font size to 44 points
        title.text_frame.paragraphs[0].font.bold = True  # Make the font bold
        # put the title little bit lower
        title.text_frame.margin_top = Inches(1)
        slide.shapes.add_picture("logot\\vuokattisportlogo.jpg", Inches(6.5), Inches(0.2), height=Inches(2.5))
        slide.shapes.add_picture("logot\\fis-main-logo-rvb.png", Inches(0.6), Inches(0.4), height=Inches(2))

        logging.info(f"Adding moving gif {combined_gif_path} to next slides")
        start_position_inches = 8   # Starting position in inches from the top of the slide
        end_position_inches = 3     # Ending position in inches from the top of the slide
        total_distance_inches = start_position_inches - end_position_inches
        slide_count = number_of_slides
        increment_per_slide = total_distance_inches / (slide_count - 1)  # Subtract 1 because you start from the first slide
        for i in range(slide_count):
            slide = prs.slides.add_slide(slide_layout)
            #top_position = Inches(8 - 0.05 * i)  # Start from the bottom and move up 0.2 inches each time
            top_position = start_position_inches - (increment_per_slide * i)
            slide.shapes.add_picture(combined_gif_path, Inches(0.8), Inches(top_position), height=Inches(4))
            slide.shapes.add_picture("logot\\vuokattisportlogo.jpg", Inches(6.5), Inches(0.2), height=Inches(2.5))
            slide.shapes.add_picture("logot\\fis-main-logo-rvb.png", Inches(0.6), Inches(0.4), height=Inches(2))
        logging.info(f"created slides for sarja {sarja}")
    logging.info(f"Saving presentation into Rising_Flags_{day_number}_same_{flag_size}.pptx")
    # Save to a user-accessible directory
    now = datetime.now()
    date_time_string = now.strftime("%Y_%m_%d_%H_%M_%S")


    final_pptx_path = os.path.join(os.getcwd(), f'This_is_temp_file_do_not_use_{date_time_string}.pptx')
    prs.save(final_pptx_path)
    return final_pptx_path, winners_acronyms

# Placeholder function for your computation that creates a PowerPoint file
def create_presentation_qt_func(csv_file_path,databack, prefix,suffix):
    data = databack
    filename = csv_file_path
    day_number = 2
    flag_size = "height" # "height" or "width"
    number_of_slides = 100
    """
    if len(sys.argv) != 5:
        print(f"usage: python3 rising.py <csv_file> <day_number> <flag_size=['height','width'> <number_of_slides>")
        exit(1)
    if int(sys.argv[4]) != 100:
        print(f"warning: number of slides is not 100, but {sys.argv[4]}")
        exit(1)
    """

    output_file,winners_acronyms = create_presentation(day_number, flag_size, 100,prefix,suffix,data)
    # current working directory folder:
    license = slides.License()
    license.set_license('logot\\Aspose.SlidesforPythonvia.NET.lic')
    if license.is_licensed():
        print("License is good!")
    else:
        print("licence failed")
        exit(1)

    with slides.Presentation(output_file) as presentation:
        # Iterate over each slide and set transitions
        for i in range(len(presentation.slides)):
            slide = presentation.slides[i]
            transition = slide.slide_show_transition

            if i % 101 == 0 or i % 101 == 100:  # First or second slide in each group
                # Set transition settings for specified slides
                transition.advance_on_click = True
                transition.advance_after = False# Setting to 0 disables advance after time
            else:
                # Set transition settings for all other slides
                transition.advance_on_click = False
                transition.advance_after = True
                transition.advance_after_time = 20  # 0.02 seconds in milliseconds


        for index, acronym in enumerate(winners_acronyms, start=1):

            slide_number = 1 + (index - 1) * 101
            slide = presentation.slides[slide_number]
            with open(f'audio_short\\{acronym}.mp3', 'rb') as in_file:
                audio_frame = slide.shapes.add_audio_frame_embedded(50,150,100,100,in_file)

                audio_frame.hide_at_showing = True # Set to hide the audio icon during the slideshow
                audio_frame.play_across_slides = True # To play across slides
                audio_frame.play_loop_mode = False # Set if the audio should loop
                audio_frame.play_mode = slides.AudioPlayModePreset.IN_CLICK_SEQUENCE # Set the audio to play automatically
                audio_frame.volume = slides.AudioVolumeMode.MEDIUM
        # Save the presentation
        now = datetime.now()
        date_time_string = now.strftime("%Y_%m_%d_%H_%M_%S")
        file_out = os.path.join(os.getcwd(),f"presentation_{date_time_string}.pptx")
        os.remove(output_file)
        presentation.save(file_out, slides.export.SaveFormat.PPTX)
        return file_out

class App(QWidget):
    def __init__(self):
        super().__init__()
        self.title = 'CSV to PPTX Converter'
        self.initUI()

    def initUI(self):
        self.setWindowTitle(self.title)
        layout = QVBoxLayout()

        self.prefixInput = QLineEdit(self)
        self.prefixInput.setPlaceholderText("Enter text before the event name")
        layout.addWidget(self.prefixInput)

        self.suffixInput = QLineEdit(self)
        self.suffixInput.setPlaceholderText("Enter text after the event name")
        layout.addWidget(self.suffixInput)

        self.openFileButton = QPushButton('Open CSV File')
        self.openFileButton.clicked.connect(self.openFileNameDialog)
        layout.addWidget(self.openFileButton)

        self.orderInput = QLineEdit(self)
        layout.addWidget(self.orderInput)


        self.createButton = QPushButton('Create Presentation in above order')
        self.createButton.clicked.connect(self.createPresentation_qt)
        self.createButton.setEnabled(False)  # Initially disabled
        layout.addWidget(self.createButton)

        self.downloadButton = QPushButton('Download Presentation')
        self.downloadButton.clicked.connect(self.downloadPresentation)
        self.downloadButton.setEnabled(False)  # Initially disabled
        layout.addWidget(self.downloadButton)

        self.filePathLabel = QLabel('No file selected')
        layout.addWidget(self.filePathLabel)

        # Loading GIF
        self.loadingLabel = QLabel()
        self.loadingMovie = QMovie("loading.gif")
        self.loadingLabel.setMovie(self.loadingMovie)
        self.loadingLabel.setAlignment(Qt.AlignmentFlag.AlignCenter)
        layout.addWidget(self.loadingLabel)

        self.data = None

        self.setLayout(layout)


    def onPresentationCreated(self, pptxPath):
        self.pptxPath = pptxPath
        self.filePathLabel.setText(f"Presentation created! You can download it now. '{self.pptxPath}'")
        self.downloadButton.setEnabled(True)
        self.createButton.setEnabled(True)

    @pyqtSlot()
    def openFileNameDialog(self):
        fileName, _ = QFileDialog.getOpenFileName(self, "Select a CSV file", "", "CSV Files (*.csv)")
        if fileName:
            self.filePath = fileName
            self.filePathLabel.setText(f"Selected file: {fileName}. Creating presentation takes like 15 seconds so hang in there:D")
            backdata,order = read_csv(fileName)
            self.data = backdata
            self.orderInput.setText(",".join(order))
            self.createButton.setEnabled(True)



    @pyqtSlot()
    def createPresentation_qt(self):
        if hasattr(self, 'filePath'):
            self.createButton.setEnabled(False)  # Disable the button to prevent multiple clicks
            self.filePathLabel.setText(f"Creating Presentation... it takes like 15 seconds")
            #self.pptxPath = create_presentation_qt_func(self.filePath, self.textInput.text())
            order = None
            try:
                order = self.orderInput.text().split(',')
                for key in order:
                    if key not in self.data.keys():
                        self.filePathLabel.setText("include all original items.")
                        return
                #new_data = {key: data[key] for key in order if key in data}
                new_data = dict()
                for key in order:
                    new_data[key] = self.data[key]
                self.data = new_data

                #data = new_data
            except:
                self.filePathLabel.setText("Order format not right. Must be separated by ',' and include all original items.")
                return
            self.thread = PresentationThread(self.filePath,self.data, self.prefixInput.text(), self.suffixInput.text())
            self.thread.finished.connect(self.onPresentationCreated)
            self.thread.start()
        else:
            self.filePathLabel.setText("Please select a CSV file first.")

    @pyqtSlot()
    def downloadPresentation(self):
        if hasattr(self, 'pptxPath'):
            # Logic to download or open the PPTX file
            # For now, just show a message
            self.filePathLabel.setText(f"Downloading...")
            self.filePathLabel.setText(f"Done! Presentation should now be in your computer as '{self.pptxPath}'")
        else:
            self.filePathLabel.setText("Please create a presentation first.")

class PresentationThread(QThread):
    finished = pyqtSignal(str)

    def __init__(self, filePath,databack, prefixInput, suffixInput):
        super().__init__()
        self.filePath = filePath
        self.databack = databack
        self.prefixInput = prefixInput
        self.suffixInput = suffixInput

    def run(self):
        try:
            # Modify the following line to include prefix and suffix in your presentation creation logic
            pptxPath = create_presentation_qt_func(self.filePath, self.databack, self.prefixInput, self.suffixInput)
            self.finished.emit(pptxPath)
        except Exception as e:
            print("Error during presentation creation:", e)
            self.finished.emit('')

if __name__ == '__main__':
    app = QApplication(sys.argv)
    ex = App()
    ex.show()
    sys.exit(app.exec())
