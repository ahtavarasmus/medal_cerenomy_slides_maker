import sys
import tempfile
from PyQt6.QtWidgets import QApplication, QWidget, QVBoxLayout, QPushButton, QFileDialog, QLabel, QLineEdit, QCheckBox
from PyQt6.QtGui import QMovie
from PyQt6.QtCore import pyqtSlot, Qt, QThread, pyqtSignal
from time import sleep
from pptx import Presentation
from pptx.util import Pt
import csv,sys,logging
from pptx.util import Inches
from datetime import datetime

from PIL import Image
import imageio

import aspose.slides as slides
import os


logging.basicConfig(level=logging.DEBUG,  # Set the lowest level to capture everything
                    format='%(asctime)s - %(levelname)s - %(message)s',
                    filename='app_log.txt',  # File where logs will be written
                    filemode='w')  # 'w' for overwrite each time; use 'a' to append


basedir = os.path.dirname(__file__)

# Define GIF paths and output path

gif_paths_same_height = {"fin": "masters_flags\\Flag_of_Finland.gif",
                 "aut": "masters_flags\\Flag_of_Austria.gif",
                 "ger": "masters_flags\\Flag_of_Germany.gif",
                 "usa": "masters_flags\\Flag_of_the_United_States.gif",
                 "lat": "masters_flags\\Flag_of_Latvia.gif",
                 "est": "masters_flags\\Flag_of_Estonia.gif",
                 "nzl": "masters_flags\\Flag_of_New_Zealand.gif",
                 "aus": "masters_flags\\Flag_of_Australia.gif",
                 "den": "masters_flags\\Flag_of_Denmark.gif",
                 "ita": "masters_flags\\Flag_of_Italy.gif",
                 "swe": "masters_flags\\Flag_of_Sweden.gif",
                 "jpn": "masters_flags\\Flag_of_Japan.gif",
                 "cze": "masters_flags\\Flag_of_the_Czech_Republic.gif",
                 "kaz": "masters_flags\\Flag_of_Kazakhstan.gif",
                 "nor": "masters_flags\\Flag_of_Norway.gif",
                 "ukr": "masters_flags\\Flag_of_Ukraine.gif",
                 "can": "masters_flags\\Flag_of_Canada.gif",
                 "isr": "masters_flags\\Flag_of_Israel.gif",
                 "svk": "masters_flags\\Flag_of_Slovakia.gif",
                 "sui": "masters_flags\\Flag_of_Switzerland.gif",
                 "gbr": "masters_flags\\Flag_of_the_United_Kingdom.gif",
                 "grl": "masters_flags\\Flag_of_Greenland.gif",
                 "pol": "masters_flags\\Flag_of_Poland.gif",
                 "esp": "masters_flags\\Flag_of_Spain.gif",
                 "fra": "masters_flags\\Flag_of_France.gif",
                 "rom": "masters_flags\\Flag_of_Romania.gif",
                 "arg": "masters_flags\\Flag_of_Argentina.gif"
                 }




def read_csv(filename):
    logging.info(f"Reading {filename}")

    try:
        with open(filename, mode='r', encoding='iso-8859-1') as file:
            data = dict()
            order = []
            reader = csv.reader(file, delimiter=';')
            for row in reader:
                if row[5] != "00:00:00,0":
                    data[row[0]] = [row[1:]] if row[0] not in data.keys() else data[row[0]] + [row[1:]]
                else:
                    print("skipping competitor:",row[2],"since they had time 00:00:00,0")

            for key in data.keys():
                order.append(key)
            return data,order
    except FileNotFoundError:
        logging.error(f"File {filename} not found")
        sys.exit(1)
    logging.info(f"Read {len(data)} rows")
    logging.info(f"Sarjat: {data.keys()}")


def combine_gif_rising_first(gif_paths, sarja):
    output_path = os.path.join(os.getcwd(),f"Combined_Flag_Podium_{sarja}_rising.gif")
    print("Combining GIFs rising")

    # Create a temporary file for the combined gif
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.gif')
    output_path = temp_file.name
    temp_file.close()  # Close the file so that other functions can use iti

    # Read individual GIFs and get their frames
    gif_frames = [imageio.mimread(gif) for gif in gif_paths]
    frames_len = len(gif_frames[0])
    extra_space = 50  # Additional horizontal space between flags

    # Get dimensions for combined gif
    combined_width = Image.open(gif_paths[0]).size[0]*3 + 150
    combined_height = int((max(Image.open(gif_paths[i]).size[1] for i in range(len(gif_paths))) + 300)*1.4)

    # Calculate the total number of frames
    #total_frames = 100 # works with duration 0.1 and normal height
    total_frames = 200

    print("still here haha")

    # Create frames for combined gif
    combined_frames = []
    last_y_pos = 0
    frame_idx = 0
    for frame_num in range(total_frames):
        logging.info(f"Combining frame {frame_num} of {total_frames}")
        # Create a new image with white background
        with Image.new('RGBA', (combined_width, combined_height), "WHITE") as new_frame:

            x_offset = 20
            print("before frames")
            frames = gif_frames[0]
            print("after frames")
            # Paste the individual frame
            frame_image = Image.fromarray(frames[frame_idx]) # Use the first frame of each GIF

            # Calculate the y-position for this frame
            y_pos = int((combined_height) * (frame_num / total_frames))

            y_offset = 20  # Move it higher up

            x_offset += frame_image.width + extra_space  # Add extra space after each flag
            new_frame.paste(frame_image, (x_offset,combined_height - y_pos + y_offset)) # PUT THE OFFSET HERE
            last_y_pos = combined_height - y_pos

            frame_idx = (frame_idx+1) % frames_len
            combined_frames.append(new_frame)

    print("saving")
    # Save frames as new GIF
    logging.info(f"Saving combined gif rising to {output_path} and last_y_pos is {last_y_pos}")
    imageio.mimsave(output_path, [frame for frame in combined_frames], format='GIF', duration=0.05, loop=0)
    logging.info(f"Saved combined gif rising to {output_path}")
    return output_path, frame_idx

def combine_gifs_rising_second_first(gif_paths, sarja):
    output_path = os.path.join(os.getcwd(),f"Combined_Flag_Podium_{sarja}_rising.gif")
    logging.info("Combining GIFs rising")

    # Create a temporary file for the combined gif
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.gif')
    output_path = temp_file.name
    temp_file.close()  # Close the file so that other functions can use iti

    # Read individual GIFs and get their frames
    gif_frames = [imageio.mimread(gif) for gif in gif_paths]
    frames_len = len(gif_frames[0])
    extra_space = 50  # Additional horizontal space between flags

    # Get dimensions for combined gif
    maximum_width = max(Image.open(gif_paths[i]).size[0] for i in range(len(gif_paths)))
    combined_width = maximum_width*3 + 150
    combined_height = int((max(Image.open(gif_paths[i]).size[1] for i in range(len(gif_paths))) + 300)*1.4)

    # Calculate the total number of frames
    #total_frames = 100 # works with duration 0.1 and normal height
    total_frames = 200


    # Create frames for combined gif
    combined_frames = []
    last_y_pos = 0
    frame_idx = 0
    print("combine_gifs_rising ",gif_paths)
    for frame_num in range(total_frames):
        # Create a new image with white background
        with Image.new('RGBA', (combined_width, combined_height), (255,255,255,255)) as new_frame:

            x_offset = 20
            y_offset = 150  # Default y-offset(for the first one
            for idx, frames in enumerate(gif_frames):
                # Paste the individual frame
                frame_image = Image.fromarray(frames[frame_idx]) # Use the first frame of each GIF

                # Calculate the y-position for this frame
                y_pos = int((combined_height) * (frame_num / total_frames))


                new_frame.paste(frame_image, (x_offset,combined_height - y_pos + y_offset)) # PUT THE OFFSET HERE
                last_y_pos = combined_height - y_pos
                x_offset += frame_image.width + extra_space  # Add extra space after each flag
                y_offset = 20  # Move it higher up

            frame_idx = (frame_idx+1) % frames_len
            combined_frames.append(new_frame)

    # Save frames as new GIF
    logging.info(f"Saving combined gif rising to {output_path} and last_y_pos is {last_y_pos}")
    imageio.mimsave(output_path, [frame for frame in combined_frames], format='GIF', duration=0.05, loop=0)
    logging.info(f"Saved combined gif rising to {output_path}")
    return output_path, frame_idx

def combine_gifs_rising(gif_paths, sarja):
    output_path = os.path.join(os.getcwd(),f"Combined_Flag_Podium_{sarja}_rising.gif")
    logging.info("Combining GIFs rising")
    assert len(gif_paths) == 3, f"only three GIFs must be provided, got {len(gif_paths)} with {gif_paths}"

    # Create a temporary file for the combined gif
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.gif')
    output_path = temp_file.name
    temp_file.close()  # Close the file so that other functions can use iti

    # Read individual GIFs and get their frames
    gif_frames = [imageio.mimread(gif) for gif in gif_paths]
    frames_len = len(gif_frames[0])
    extra_space = 50  # Additional horizontal space between flags

    # Get dimensions for combined gif
    combined_width = sum(Image.open(gif_paths[i]).size[0] for i in range(len(gif_paths))) + 150
    combined_height = int((max(Image.open(gif_paths[i]).size[1] for i in range(len(gif_paths))) + 300)*1.4)

    # Calculate the total number of frames
    #total_frames = 100 # works with duration 0.1 and normal height
    total_frames = 200


    # Create frames for combined gif
    combined_frames = []
    last_y_pos = 0
    frame_idx = 0
    print("combine_gifs_rising ",gif_paths)
    for frame_num in range(total_frames):
        # Create a new image with white background
        with Image.new('RGBA', (combined_width, combined_height), (255,255,255,255)) as new_frame:

            x_offset = 20
            for idx, frames in enumerate(gif_frames):
                # Paste the individual frame
                frame_image = Image.fromarray(frames[frame_idx]) # Use the first frame of each GIF

                # Calculate the y-position for this frame
                y_pos = int((combined_height) * (frame_num / total_frames))

                y_offset = 150  # Default y-offset(for the first one
                if idx == 1:  # If it's the middle (winner) flag
                    y_offset = 20  # Move it higher up
                elif idx == 2:
                    y_offset = 250

                new_frame.paste(frame_image, (x_offset,combined_height - y_pos + y_offset)) # PUT THE OFFSET HERE
                last_y_pos = combined_height - y_pos
                x_offset += frame_image.width + extra_space  # Add extra space after each flag

            frame_idx = (frame_idx+1) % frames_len
            combined_frames.append(new_frame)

    # Save frames as new GIF
    logging.info(f"Saving combined gif rising to {output_path} and last_y_pos is {last_y_pos}")
    imageio.mimsave(output_path, [frame for frame in combined_frames], format='GIF', duration=0.05, loop=0)
    logging.info(f"Saved combined gif rising to {output_path}")
    return output_path, frame_idx

def combine_gif_first(gif_paths, sarja, frame_idx):
    output_path = os.path.join(os.getcwd(),"Combined_Flag_Podium_{sarja}.gif")
    logging.info("Combining GIFs")

     # Create a temporary file for the combined gif
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.gif')
    output_path = temp_file.name
    temp_file.close()  # Close the file so that other functions can use iti

    # Read individual GIFs and get their frames
    gif_frames = [imageio.mimread(gif) for gif in gif_paths]
    frames_len = len(gif_frames[0])
    extra_space = 50  # Additional horizontal space between flags

    # Get dimensions for combined gif
    combined_width = Image.open(gif_paths[0]).size[0]*3 + 150
    combined_height = max(Image.open(gif_paths[i]).size[1] for i in range(len(gif_paths))) + 300

    # Create frames for combined gif
    combined_frames = []
    i = frame_idx
    for _ in range(frames_len):
        logging.info(f"Combining frame {i} of {min(len(frames) for frames in gif_frames)}")
        # Create a new image with white background
        with Image.new('RGBA', (combined_width, combined_height), "WHITE") as new_frame:

            x_offset = 20
            frames = gif_frames[0]
            # Paste the individual frame
            frame_image = Image.fromarray(frames[i])

            y_offset = 20  # Move it higher up

            x_offset += frame_image.width + extra_space  # Add extra space after each flag
            new_frame.paste(frame_image, (x_offset, y_offset+5))

            combined_frames.append(new_frame)
            i = (i+1) % frames_len

    # Save frames as new GIF
    logging.info(f"Saving combined gif to {output_path}")
    imageio.mimsave(output_path, [frame for frame in combined_frames], format='GIF', duration=0.2, loop=0)
    logging.info(f"Saved combined gif to {output_path}")
    return output_path

def combine_gifs_second_first(gif_paths, sarja, frame_idx):
    output_path = os.path.join(os.getcwd(),"Combined_Flag_Podium_{sarja}.gif")
    logging.info("Combining GIFs")

     # Create a temporary file for the combined gif
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.gif')
    output_path = temp_file.name
    temp_file.close()  # Close the file so that other functions can use iti

    # Read individual GIFs and get their frames
    gif_frames = [imageio.mimread(gif) for gif in gif_paths]
    frames_len = len(gif_frames[0])
    extra_space = 50  # Additional horizontal space between flags

    # Get dimensions for combined gif
    maximum_width = max(Image.open(gif_paths[i]).size[0] for i in range(len(gif_paths)))
    combined_width = maximum_width*3 + 150
    combined_height = max(Image.open(gif_paths[i]).size[1] for i in range(len(gif_paths))) + 300

    # Create frames for combined gif
    combined_frames = []
    i = frame_idx
    print("combine_gifs ",gif_paths)
    for _ in range(frames_len):
        # Create a new image with white background
        with Image.new('RGBA', (combined_width, combined_height), (255,255,255,255)) as new_frame:

            x_offset = 20
            y_offset = 150  # Default y-offset
            for idx, frames in enumerate(gif_frames):
                # Paste the individual frame
                frame_image = Image.fromarray(frames[i])


                new_frame.paste(frame_image, (x_offset, y_offset+5))
                x_offset += frame_image.width + extra_space  # Add extra space after each flag
                y_offset = 20  # Move it higher up

            combined_frames.append(new_frame)
            i = (i+1) % frames_len

    # Save frames as new GIF
    logging.info(f"Saving combined gif to {output_path}")
    imageio.mimsave(output_path, [frame for frame in combined_frames], format='GIF', duration=0.2, loop=0)
    logging.info(f"Saved combined gif to {output_path}")
    return output_path



def combine_gifs(gif_paths, sarja, frame_idx):
    output_path = os.path.join(os.getcwd(),"Combined_Flag_Podium_{sarja}.gif")
    logging.info("Combining GIFs")
    assert len(gif_paths) == 3, f"only three GIFs must be provided, got {len(gif_paths)} with {gif_paths}"

     # Create a temporary file for the combined gif
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.gif')
    output_path = temp_file.name
    temp_file.close()  # Close the file so that other functions can use iti

    # Read individual GIFs and get their frames
    gif_frames = [imageio.mimread(gif) for gif in gif_paths]
    frames_len = len(gif_frames[0])
    extra_space = 50  # Additional horizontal space between flags

    # Get dimensions for combined gif
    combined_width = sum(Image.open(gif_paths[i]).size[0] for i in range(len(gif_paths))) + 150
    combined_height = max(Image.open(gif_paths[i]).size[1] for i in range(len(gif_paths))) + 300

    # Create frames for combined gif
    combined_frames = []
    i = frame_idx
    print("combine_gifs ",gif_paths)
    for _ in range(frames_len):
        # Create a new image with white background
        with Image.new('RGBA', (combined_width, combined_height), (255,255,255,255)) as new_frame:

            x_offset = 20
            for idx, frames in enumerate(gif_frames):
                # Paste the individual frame
                frame_image = Image.fromarray(frames[i])

                y_offset = 150  # Default y-offset
                if idx == 1:  # If it's the middle (winner) flag
                    y_offset = 20  # Move it higher up
                elif idx == 2:
                    y_offset = 250

                new_frame.paste(frame_image, (x_offset, y_offset+5))
                x_offset += frame_image.width + extra_space  # Add extra space after each flag

            combined_frames.append(new_frame)
            i = (i+1) % frames_len

    # Save frames as new GIF
    logging.info(f"Saving combined gif to {output_path}")
    imageio.mimsave(output_path, [frame for frame in combined_frames], format='GIF', duration=0.2, loop=0)
    logging.info(f"Saved combined gif to {output_path}")
    return output_path




basedir = os.path.dirname(__file__)

def abs_path(relative_path):
    """ Get the absolute path to the resource """

    return os.path.join(basedir, relative_path)

def create_presentation(day_number, number_of_slides,prefix,suffix,data,infoBox,namesChecked,staticChecked):
    logging.info("Creating presentation")
    prs = Presentation()
    slide_layout = prs.slide_layouts[0] # blank slide

    winners_acronyms = []
    for sarja in data.keys():
        # creating 100 slides for each sarja to create animation
        # sorting based on tulos
        logging.info(f"sarja {sarja}, data[sarja]: ")
        infoBox.setText(f"making slides for sarja: {sarja}")
        for row in data[sarja]:
            logging.info(row)
        #logging.info(f"Sorting sarja {sarja}, data[sarja]: ")
        #sorted(data[sarja], key=lambda x: x[4])

        if len(data[sarja]) == 1:
            print(f"one row for sarja:{data[sarja]}")
            first_country = data[sarja][0][3].lower()
            print("after")
            first_name = data[sarja][0][1]
            print("after moi")
            first = gif_paths_same_height[first_country]
            print("after moi")
            second_country = None
            third_country = None
        elif len(data[sarja]) == 2:
            print(f"only two on the podium")
            first_country = data[sarja][0][3].lower()
            first_name = data[sarja][0][1]
            first = gif_paths_same_height[first_country]

            second_country = data[sarja][1][3].lower()
            second_name = data[sarja][1][1]
            second = gif_paths_same_height[second_country]

            third_country = None

        else:
            first_country = data[sarja][0][3].lower()
            first_name = data[sarja][0][1]
            second_country = data[sarja][1][3].lower()
            second_name = data[sarja][1][1]
            third_country = data[sarja][2][3].lower()
            third_name = data[sarja][2][1]
            if first_country in gif_paths_same_height.keys() and second_country in gif_paths_same_height.keys() and third_country in gif_paths_same_height.keys():
                first = gif_paths_same_height[first_country]
                second = gif_paths_same_height[second_country]
                third = gif_paths_same_height[third_country]
            else:
                logging.error(f"didn't find country for sarja {sarja} with countries {first_country}, {second_country}, {third_country}")
                continue

            if not first or not second or not third:
                logging.error(f"didn't find gif for sarja {sarja} with countries {first_country}, {second_country}, {third_country}")
                continue
            print(f"Sarja: {sarja}, 1: {first}, 2: {second}, 3: {third}")

        winners_acronyms.append(first_country.upper())
        logging.info(f"Making slide for showing sarja {sarja}")

        slide = prs.slides.add_slide(slide_layout)
        #left = top = Inches(0)
        #slide.shapes.add_picture('logot\\background.jpg',left,top,width=Inches(10),height=Inches(7.5))
        # adding stationary large title text to the middle
        title = slide.shapes.title
        title.text = prefix+sarja+suffix
        print("here?")
        title.text_frame.paragraphs[0].font.size = Pt(44)  # Set font size to 44 points
        title.text_frame.paragraphs[0].font.bold = True  # Make the font bold
        # put the title little bit lower
        title.text_frame.margin_top = Inches(1)

        vuokattisport_left = Inches(7.5)
        vuokattisport_top = Inches(0.3)
        vuokattisport_height = Inches(1.8)
        masters_left = Inches(0.5)
        masters_top = Inches(0.3)
        masters_height = Inches(1.5)
        slide.shapes.add_picture("logot\\Vuokatti_sport_logo.png", vuokattisport_left, vuokattisport_top, height=vuokattisport_height)
        slide.shapes.add_picture("logot\\masterslogo.png", masters_left, masters_top, height=masters_height)

        print("here??")


        slide = prs.slides.add_slide(slide_layout)
        #left = top = Inches(0)
        #slide.shapes.add_picture('logot\\background.jpg',left,top,width=Inches(10),height=Inches(7.5))
        # adding stationary large title text to the middle
        title = slide.shapes.title
        title.text = prefix+sarja+suffix
        title.text_frame.paragraphs[0].font.size = Pt(44)  # Set font size to 44 points
        title.text_frame.paragraphs[0].font.bold = True  # Make the font bold
        # put the title little bit lower
        title.text_frame.margin_top = Inches(1)

        print("here???")

        if namesChecked:
            left = Inches(3.3)
            top = Inches(4)
            width = Inches(5)
            height = Inches(2)
            txBox = slide.shapes.add_textbox(left,top,width,height)
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = f"1. {first_name}"
            tf.paragraphs[0].font.size = Pt(33)
            tf.paragraphs[1].font.size = Pt(33)
            print("here????")
            if second_country:
                p = tf.add_paragraph()
                p.text = f"2. {second_name}"
                tf.paragraphs[2].font.size = Pt(33)
            if third_country:
                p = tf.add_paragraph()
                p.text = f"3. {third_name}"
                tf.paragraphs[3].font.size = Pt(33)

        slide.shapes.add_picture("logot\\Vuokatti_sport_logo.png", vuokattisport_left, vuokattisport_top, height=vuokattisport_height)
        slide.shapes.add_picture("logot\\masterslogo.png", masters_left, masters_top, height=masters_height)

        flags_left = Inches(1.3) # was 2.3(smaller)
        flags_top = Inches(2.5) # was 3(smaller)
        flags_width = Inches(8) # was 6(smaller)

        frame_idx = 0
        if not staticChecked:
            if not second_country: # only winners flag
                print("here moihei")
                combined_gif_path_rising,frame_idx = combine_gif_rising_first([first],sarja)
                print("sitll here")
            elif not third_country: # only two first ones
                print("only two competitors in one sarja")
                try:
                    combined_gif_path_rising,frame_idx = combine_gifs_rising_second_first([second,first],sarja)
                except:
                    print("FAILED RISING IN SARJA",sarja)
                    combined_gif_path_rising = None
            else:
                try:
                    combined_gif_path_rising,frame_idx = combine_gifs_rising([second, first, third], sarja)
                except:
                    print("FAILED RISING IN SARJA",sarja)
                    combined_gif_path_rising = None
            if combined_gif_path_rising:
                # rising(flags move from bottom to their position) gif slide
                slide = prs.slides.add_slide(slide_layout)
                #slide.shapes.add_picture('logot\\background.jpg',left,top,width=Inches(10),height=Inches(7.5))
                slide.shapes.add_picture(combined_gif_path_rising, flags_left,flags_top,width=flags_width)
                slide.shapes.add_picture("logot\\Vuokatti_sport_logo.png", vuokattisport_left, vuokattisport_top, height=vuokattisport_height)
                slide.shapes.add_picture("logot\\masterslogo.png", masters_left, masters_top, height=masters_height)
            else:
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.add_picture("logot\\Vuokatti_sport_logo.png", vuokattisport_left, vuokattisport_top, height=vuokattisport_height)
                slide.shapes.add_picture("logot\\masterslogo.png", masters_left, masters_top, height=masters_height)



        if not second_country: # only winners flag
            print("only one competitor in sarja")
            try:
                combined_gif_path = combine_gif_first([first],sarja,frame_idx)
            except:
                print("FAILED STATIC IN SARJA",sarja)
                combined_gif_path = None
        elif not third_country: # only two first ones
            print("only two comps in one sarja")
            try:
                combined_gif_path = combine_gifs_second_first([second,first],sarja,frame_idx)
            except:
                print("FAILED STATIC IN SARJA",sarja)
                combined_gif_path = None
        else:
            try:
                combined_gif_path = combine_gifs([second, first, third], sarja, frame_idx)
            except:
                print("FAILED STATIC IN SARJA",sarja)
                combined_gif_path = None
        if combined_gif_path:
            # 'static'(actually the flags are moving in place) gif slide
            slide = prs.slides.add_slide(slide_layout)
            #slide.shapes.add_picture('logot\\background.jpg',left,top,width=Inches(10),height=Inches(7.5))
            slide.shapes.add_picture(combined_gif_path, flags_left,flags_top,width=flags_width)
            slide.shapes.add_picture("logot\\Vuokatti_sport_logo.png", vuokattisport_left, vuokattisport_top, height=vuokattisport_height)
            slide.shapes.add_picture("logot\\masterslogo.png", masters_left, masters_top, height=masters_height)
        else:
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.add_picture("logot\\Vuokatti_sport_logo.png", vuokattisport_left, vuokattisport_top, height=vuokattisport_height)
            slide.shapes.add_picture("logot\\masterslogo.png", masters_left, masters_top, height=masters_height)



        logging.info(f"created slides for sarja {sarja}")

    # Save to a user-accessible directory
    now = datetime.now()
    date_time_string = now.strftime("%Y_%m_%d_%H_%M_%S")


    final_pptx_path = os.path.join(os.getcwd(), f'DO_NOT_OPEN_This_file_is_temporary{date_time_string}.pptx')
    prs.save(final_pptx_path)
    return final_pptx_path, winners_acronyms

# Placeholder function for your computation that creates a PowerPoint file
def create_presentation_qt_func(csv_file_path,databack, prefix,suffix, infoBox,namesChecked,staticChecked):
    data = databack
    filename = csv_file_path
    day_number = 2
    number_of_slides = 100
    """
    if len(sys.argv) != 5:
        print(f"usage: python3 rising.py <csv_file> <day_number> <flag_size=['height','width'> <number_of_slides>")
        exit(1)
    if int(sys.argv[4]) != 100:
        print(f"warning: number of slides is not 100, but {sys.argv[4]}")
        exit(1)
    """

    output_file,winners_acronyms = create_presentation(day_number, 100,prefix,suffix,data, infoBox, namesChecked, staticChecked)
    # current working directory folder:
    license = slides.License()
    license.set_license('logot\\Aspose.SlidesforPythonvia.NET.lic')
    if license.is_licensed():
        print("License is good!")
    else:
        print("licence failed")
        exit(1)

    infoBox.setText(f"adding audio and modifying transitions to the slides")
    with slides.Presentation(output_file) as presentation:

        # Iterate over each slide and set transitions
        if not staticChecked:
            rising_slide_index = 2
            for i in range(len(presentation.slides)):
                slide = presentation.slides[i]
                transition = slide.slide_show_transition
                if i == rising_slide_index:
                    # Set transition settings for rising flag slide
                    transition.advance_on_click = False
                    transition.advance_after = True
                    transition.advance_after_time = 19400  # 19.35 seconds in milliseconds
                    rising_slide_index += 4 # setting for the next sarja
                else:
                    # Set transition settings for other slides(sarja slide and static flag slide)
                    transition.advance_on_click = True
                    transition.advance_after = False


        music_slide_index = 2
        for acronym in winners_acronyms:
            slide = presentation.slides[music_slide_index]
            with open(f'audio_short\\{acronym}.mp3', 'rb') as in_file:
                audio_frame = slide.shapes.add_audio_frame_embedded(50,150,100,100,in_file)

                audio_frame.hide_at_showing = True # Set to hide the audio icon during the slideshow
                audio_frame.play_across_slides = True # To play across slides
                audio_frame.play_loop_mode = False # Set if the audio should loop
                audio_frame.play_mode = slides.AudioPlayModePreset.AUTO # Set the audio to play automatically
                audio_frame.volume = slides.AudioVolumeMode.MEDIUM
            if staticChecked:
                music_slide_index += 3 # set for the next static slide
            else:
                music_slide_index += 4 # set for the next rising slide
        # Save the presentation
        now = datetime.now()
        date_time_string = now.strftime("%Y_%m_%d_%H_%M_%S")
        if staticChecked:
            file_out = os.path.join(os.getcwd(),f"BACKUP_WITH_STATIC_FLAGS_{date_time_string}.pptx")
        else:
            file_out = os.path.join(os.getcwd(),f"Ceremony_With_Rising_Flags_{date_time_string}.pptx")
        infoBox.setText(f"done!  '{file_out}'")
        os.remove(output_file)
        presentation.save(file_out, slides.export.SaveFormat.PPTX)
        return file_out

class App(QWidget):
    def __init__(self):
        super().__init__()
        self.title = 'CSV to PPTX Converter'
        self.initUI()

    def initUI(self):
        self.setWindowTitle(self.title)
        layout = QVBoxLayout()

        self.prefixInput = QLineEdit(self)
        self.prefixInput.setPlaceholderText("Enter text before the event name")
        layout.addWidget(self.prefixInput)

        self.suffixInput = QLineEdit(self)
        self.suffixInput.setPlaceholderText("Enter text after the event name")
        layout.addWidget(self.suffixInput)

        self.openFileButton = QPushButton('Open CSV File')
        self.openFileButton.clicked.connect(self.openFileNameDialog)
        layout.addWidget(self.openFileButton)

        self.orderInput = QLineEdit(self)
        layout.addWidget(self.orderInput)

        self.namesCheckBox = QCheckBox('Names under sarja',self)
        layout.addWidget(self.namesCheckBox)

        self.staticCheckBox = QCheckBox('Static flags',self)
        layout.addWidget(self.staticCheckBox)


        self.createButton = QPushButton('Create Presentation in above order')
        self.createButton.clicked.connect(self.createPresentation_qt)
        self.createButton.setEnabled(False)  # Initially disabled
        layout.addWidget(self.createButton)


        self.infoBox = QLabel('No file selected')
        layout.addWidget(self.infoBox)

        # Loading GIF
        self.loadingLabel = QLabel()
        self.loadingMovie = QMovie("loading.gif")
        self.loadingLabel.setMovie(self.loadingMovie)
        self.loadingLabel.setAlignment(Qt.AlignmentFlag.AlignCenter)
        layout.addWidget(self.loadingLabel)

        self.data = None

        self.setLayout(layout)


    def onPresentationCreated(self, pptxPath):
        self.pptxPath = pptxPath
        self.infoBox.setText(f"Presentation created!  '{self.pptxPath}'")
        self.namesCheckBox.setEnabled(True)
        self.staticCheckBox.setEnabled(True)
        self.createButton.setEnabled(True)
        self.orderInput.setEnabled(True)
        self.suffixInput.setEnabled(True)
        self.prefixInput.setEnabled(True)
        self.openFileButton.setEnabled(True)


    @pyqtSlot()
    def openFileNameDialog(self):
        fileName, _ = QFileDialog.getOpenFileName(self, "Select a CSV file", "", "CSV Files (*.csv)")
        if fileName:
            self.filePath = fileName
            self.infoBox.setText(f"Selected file: {fileName}. Creating presentation takes like 15 seconds so hang in there:D")
            backdata,order = read_csv(fileName)
            self.data = backdata
            self.orderInput.setText(",".join(order))
            self.createButton.setEnabled(True)



    @pyqtSlot()
    def createPresentation_qt(self):
        if hasattr(self, 'filePath'):
            self.createButton.setEnabled(False)  # Disable the button to prevent multiple clicks
            self.namesCheckBox.setEnabled(False)
            self.orderInput.setEnabled(False)
            self.suffixInput.setEnabled(False)
            self.prefixInput.setEnabled(False)
            self.openFileButton.setEnabled(False)
            self.staticCheckBox.setEnabled(False)
            self.infoBox.setText(f"Creating Presentation... takes like 15 seconds")
            #self.pptxPath = create_presentation_qt_func(self.filePath, self.textInput.text())
            order = None
            try:
                order = self.orderInput.text().split(',')
                for key in order:
                    if key not in self.data.keys():
                        self.infoBox.setText("include all original items.")
                        return
                #new_data = {key: data[key] for key in order if key in data}
                new_data = dict()
                for key in order:
                    new_data[key] = self.data[key]
                self.data = new_data

                #data = new_data
            except:
                self.infoBox.setText("Order format not right. Must be separated by ',' and include all original items.")
                return
            self.thread = PresentationThread(self.filePath,self.data, self.infoBox, self.prefixInput.text(), self.suffixInput.text(), self.namesCheckBox.isChecked(), self.staticCheckBox.isChecked())
            self.thread.finished.connect(self.onPresentationCreated)
            self.thread.start()
        else:
            self.infoBox.setText("Please select a CSV file first.")


class PresentationThread(QThread):
    finished = pyqtSignal(str)

    def __init__(self, filePath,databack, infoBox,prefixInput, suffixInput, namesChecked,staticChecked):
        super().__init__()
        self.filePath = filePath
        self.databack = databack
        self.infoBox = infoBox
        self.prefixInput = prefixInput
        self.suffixInput = suffixInput
        self.namesChecked = namesChecked
        self.staticChecked = staticChecked

    def run(self):
        try:
            # Modify the following line to include prefix and suffix in your presentation creation logic
            pptxPath = create_presentation_qt_func(self.filePath, self.databack, self.prefixInput, self.suffixInput, self.infoBox, self.namesChecked, self.staticChecked)
            self.finished.emit(pptxPath)
        except Exception as e:
            print("Error during presentation creation:", e)
            self.finished.emit('')

if __name__ == '__main__':
    app = QApplication(sys.argv)
    ex = App()
    ex.show()
    sys.exit(app.exec())
